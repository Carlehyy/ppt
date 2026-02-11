#!/usr/bin/env python3
"""
Content Parser Module — 四层递进内容解析架构
Layer 1: 物理层解析 — 格式 → 原始文本 + 元数据
Layer 2: 语义层解析 — 原始文本 → 结构化语义单元（LLM驱动）
Layer 3: 关联层分析 — 多文档 → 关系图谱（LLM驱动）
Layer 4: 汇报框架映射 — 语义单元 → PPT大纲雏形（LLM驱动）
"""

import os
import re
import json
from typing import List, Dict, Any, Optional
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
try:
    from .utils import detect_file_type, format_source_info, save_json
except ImportError:
    from utils import detect_file_type, format_source_info, save_json


class ContentParser:
    """四层递进内容解析器"""

    def __init__(self, llm_client=None):
        self.llm_client = llm_client

    # ------------------------------------------------------------------ #
    #  公共入口
    # ------------------------------------------------------------------ #
    def parse(self, file_paths: List[str]) -> Dict[str, Any]:
        """
        执行完整的四层递进解析流程

        Returns:
            {
                "documents": [...],          # 物理层产出
                "merged_content": str,        # 合并后的原始文本
                "document_profiles": [...],   # 语义层 — 文档画像
                "semantic_units": [...],      # 语义层 — 信息要素
                "document_relationships": {}, # 关联层 — 多文档关系
                "framework_mapping": {},      # 框架层 — 汇报框架映射
                "information_gaps": [...],    # 框架层 — 信息缺口
                "summary": str
            }
        """
        # ---- Layer 1: 物理层解析 ---- #
        print("  [Layer 1] 物理层解析...")
        documents = []
        for fp in file_paths:
            print(f"    解析文件: {os.path.basename(fp)}")
            ext = os.path.splitext(fp)[1].lower()
            parser_map = {
                ".docx": self._parse_docx,
                ".pdf": self._parse_pdf,
                ".pptx": self._parse_pptx,
                ".txt": self._parse_txt,
                ".md": self._parse_txt,
            }
            parser_fn = parser_map.get(ext)
            if parser_fn is None:
                print(f"    ⚠ 不支持的文件格式: {ext}，跳过")
                continue
            content = parser_fn(fp)
            documents.append({
                "file_path": fp,
                "file_name": os.path.basename(fp),
                "file_type": ext,
                "content": content,
            })
        print(f"  [Layer 1] ✓ 完成，共解析 {len(documents)} 个文件")

        merged_content = self._merge_documents(documents)

        # ---- Layer 2: 语义层解析（LLM 驱动） ---- #
        document_profiles = []
        semantic_units = []
        if self.llm_client:
            print("  [Layer 2] 语义层解析...")
            # 两遍处理策略：第一遍快速通读建立文档画像
            document_profiles = self._build_document_profiles(documents)
            # 第二遍精细提取语义单元
            semantic_units = self._extract_semantic_units(documents, document_profiles)
            # 数据敏感信息专项处理：正则交叉验证
            semantic_units = self._cross_validate_data(semantic_units, merged_content)
            print(f"  [Layer 2] ✓ 完成，提取 {len(semantic_units)} 个语义单元")
        else:
            print("  [Layer 2] ⚠ 未配置LLM，跳过语义层解析")

        # ---- Layer 3: 关联层分析（多文档） ---- #
        document_relationships = {}
        if self.llm_client and len(documents) > 1:
            print("  [Layer 3] 关联层分析...")
            document_relationships = self._analyze_relationships(document_profiles, semantic_units)
            print("  [Layer 3] ✓ 完成")
        elif len(documents) == 1:
            document_relationships = {
                "type": "single_document",
                "overall_narrative": document_profiles[0].get("core_theme", "") if document_profiles else "",
            }
            print("  [Layer 3] ✓ 单文档，跳过关联分析")

        # ---- Layer 4: 汇报框架映射 ---- #
        framework_mapping = {}
        information_gaps = []
        if self.llm_client:
            print("  [Layer 4] 汇报框架映射...")
            framework_mapping, information_gaps = self._map_to_framework(
                document_profiles, semantic_units, document_relationships
            )
            print(f"  [Layer 4] ✓ 完成，发现 {len(information_gaps)} 个信息缺口")

        return {
            "documents": documents,
            "merged_content": merged_content,
            "document_profiles": document_profiles,
            "semantic_units": semantic_units,
            "document_relationships": document_relationships,
            "framework_mapping": framework_mapping,
            "information_gaps": information_gaps,
            "summary": self._generate_summary(documents, semantic_units),
        }

    # ================================================================== #
    #  Layer 1 — 物理层解析
    # ================================================================== #
    def _parse_docx(self, file_path: str) -> Dict[str, Any]:
        """解析 Word 文档，保留标题层级、段落样式、表格、图片元数据"""
        doc = Document(file_path)

        # --- 段落 --- #
        paragraphs = []
        current_section = None
        for i, para in enumerate(doc.paragraphs):
            text = para.text.strip()
            if not text:
                continue
            style_name = para.style.name if para.style else ""
            is_heading = style_name.startswith("Heading")
            heading_level = int(style_name.replace("Heading ", "")) if is_heading and style_name[-1].isdigit() else 0

            # 检测加粗/高亮
            has_bold = any(run.bold for run in para.runs if run.bold)
            has_highlight = any(run.font.highlight_color is not None for run in para.runs)

            if is_heading:
                current_section = text

            paragraphs.append({
                "index": i,
                "text": text,
                "style": style_name,
                "is_heading": is_heading,
                "heading_level": heading_level,
                "section": current_section,
                "has_bold": has_bold,
                "has_highlight": has_highlight,
                "source": format_source_info(file_path, page_num=None),
            })

        # --- 表格 --- #
        tables = []
        for i, table in enumerate(doc.tables):
            rows_data = []
            for row in table.rows:
                rows_data.append([cell.text.strip() for cell in row.cells])
            tables.append({
                "index": i,
                "data": rows_data,
                "rows": len(rows_data),
                "cols": len(rows_data[0]) if rows_data else 0,
                "source": format_source_info(file_path),
            })

        # --- 图片计数 --- #
        image_count = 0
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                image_count += 1

        return {
            "paragraphs": paragraphs,
            "tables": tables,
            "image_count": image_count,
            "total_paragraphs": len(paragraphs),
            "total_tables": len(tables),
            "sections": self._extract_sections_from_paragraphs(paragraphs),
        }

    def _parse_pdf(self, file_path: str) -> Dict[str, Any]:
        """解析 PDF 文档"""
        reader = PdfReader(file_path)
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                pages.append({
                    "page_num": i + 1,
                    "text": text.strip(),
                    "source": format_source_info(file_path, i + 1),
                })
        return {"pages": pages, "total_pages": len(pages)}

    def _parse_pptx(self, file_path: str) -> Dict[str, Any]:
        """解析 PPT 文档，提取文本、备注、版式类型"""
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides):
            slide_info = {
                "slide_num": i + 1,
                "title": "",
                "texts": [],
                "notes": "",
                "layout_name": slide.slide_layout.name if slide.slide_layout else "",
                "source": format_source_info(file_path, i + 1),
            }
            if slide.shapes.title:
                slide_info["title"] = slide.shapes.title.text
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape != slide.shapes.title:
                    slide_info["texts"].append(shape.text)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                slide_info["notes"] = slide.notes_slide.notes_text_frame.text
            slides.append(slide_info)
        return {"slides": slides, "total_slides": len(slides)}

    def _parse_txt(self, file_path: str) -> Dict[str, Any]:
        """解析纯文本 / Markdown 文件"""
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
        lines = [l for l in text.split("\n") if l.strip()]
        return {"text": text, "lines": lines, "total_lines": len(lines)}

    # ================================================================== #
    #  Layer 2 — 语义层解析
    # ================================================================== #
    def _build_document_profiles(self, documents: List[Dict]) -> List[Dict]:
        """第一遍快速通读：为每个文档建立画像"""
        profiles = []
        for doc in documents:
            text_preview = self._get_text_preview(doc, max_chars=3000)
            prompt = (
                "你是专业的文档分析师。请对以下文档进行快速画像分析，输出JSON：\n"
                "{\n"
                '  "doc_type": "工作总结/项目报告/数据分析/方案提案/会议纪要/其他",\n'
                '  "time_range": "涉及的时间范围",\n'
                '  "core_theme": "一句话概括核心主题",\n'
                '  "info_nature": "成果展示型/问题分析型/规划建议型/综合型",\n'
                '  "key_topics": ["主题1", "主题2", ...],\n'
                '  "data_richness": "丰富/一般/匮乏",\n'
                '  "estimated_importance": "核心素材/辅助素材/参考素材"\n'
                "}\n\n"
                f"文件名: {doc['file_name']}\n"
                f"内容预览:\n{text_preview}"
            )
            try:
                result = self.llm_client.call_llm(prompt, response_json=True)
                result["file_name"] = doc["file_name"]
                profiles.append(result)
            except Exception as e:
                print(f"    ⚠ 文档画像失败({doc['file_name']}): {e}")
                profiles.append({"file_name": doc["file_name"], "doc_type": "未知", "core_theme": ""})
        return profiles

    def _extract_semantic_units(self, documents: List[Dict], profiles: List[Dict]) -> List[Dict]:
        """第二遍精细提取：按Chunk策略逐段提取语义单元"""
        all_units = []
        for doc, profile in zip(documents, profiles):
            chunks = self._chunk_document(doc)
            context_hint = json.dumps(profile, ensure_ascii=False)
            for chunk in chunks:
                prompt = (
                    "你是专业的信息提取专家。请从以下文档片段中提取结构化语义单元。\n\n"
                    f"文档画像: {context_hint}\n\n"
                    "提取规则:\n"
                    "1. 信息要素分类: background(背景/目标), achievement(关键成果), "
                    "method(方法/过程), problem(问题/风险), plan(下一步计划), "
                    "data(关键数据), conclusion(关键结论)\n"
                    "2. 颗粒度标注: must_show(必须呈现), should_show(建议呈现), optional(可选补充)\n"
                    "3. 每个单元必须标注来源(source)和置信度(confidence: high/medium/low)\n"
                    "4. 特别关注数字、百分比、金额、日期等数据敏感信息\n\n"
                    "输出JSON:\n"
                    '{"units": [\n'
                    '  {"type": "achievement", "content": "...", "granularity": "must_show", '
                    '"source": "文件名:章节/段落", "confidence": "high", '
                    '"key_data": ["30%", "1200万"], "keywords": ["效率", "成本"]}\n'
                    "]}\n\n"
                    f"文档片段 (来自 {doc['file_name']}):\n{chunk['text']}"
                )
                try:
                    result = self.llm_client.call_llm(prompt, response_json=True)
                    units = result.get("units", [])
                    # 补充来源信息
                    for u in units:
                        u.setdefault("source_file", doc["file_name"])
                        u.setdefault("chunk_section", chunk.get("section", ""))
                    all_units.extend(units)
                except Exception as e:
                    print(f"    ⚠ 语义提取失败({doc['file_name']}, chunk): {e}")
        return all_units

    def _cross_validate_data(self, units: List[Dict], merged_text: str) -> List[Dict]:
        """数据敏感信息专项处理：正则交叉验证"""
        # 从原文中提取所有数字模式
        number_patterns = re.findall(
            r'[\d,]+\.?\d*\s*[%％万亿元美元千百十]|'
            r'\d{4}[-/年]\d{1,2}[-/月]\d{0,2}日?|'
            r'\d+\.?\d*\s*(?:万|亿|千|百|元|美元|%|％)',
            merged_text
        )
        original_numbers = set(n.strip() for n in number_patterns)

        for unit in units:
            key_data = unit.get("key_data", [])
            if not key_data:
                continue
            validated = []
            for d in key_data:
                # 检查该数据是否能在原文中找到
                found = d in merged_text or any(d in n for n in original_numbers)
                validated.append({"value": d, "verified": found})
            unit["data_validation"] = validated
            # 如果有未验证的数据，降低置信度
            if any(not v["verified"] for v in validated):
                if unit.get("confidence") == "high":
                    unit["confidence"] = "medium"
                    unit["validation_warning"] = "部分数据未在原文中找到精确匹配，请用户核实"
        return units

    # ================================================================== #
    #  Layer 3 — 关联层分析
    # ================================================================== #
    def _analyze_relationships(self, profiles: List[Dict], units: List[Dict]) -> Dict:
        """分析多文档之间的关系"""
        profiles_str = json.dumps(profiles, ensure_ascii=False)
        # 统计各文档的语义单元分布
        unit_distribution = {}
        for u in units:
            fname = u.get("source_file", "unknown")
            utype = u.get("type", "other")
            unit_distribution.setdefault(fname, {})
            unit_distribution[fname][utype] = unit_distribution[fname].get(utype, 0) + 1

        prompt = (
            "你是专业的多文档分析师。请分析以下多个文档之间的关系。\n\n"
            f"文档画像:\n{profiles_str}\n\n"
            f"语义单元分布:\n{json.dumps(unit_distribution, ensure_ascii=False)}\n\n"
            "请输出JSON:\n"
            "{\n"
            '  "relationships": [\n'
            '    {"doc_a": "文件A", "doc_b": "文件B", "relation_type": "时间关系/主题关系/层级关系/补充关系", "description": "..."}\n'
            "  ],\n"
            '  "overall_narrative": "一句话描述所有文档共同讲述的故事",\n'
            '  "recommended_storyline": "建议的叙事主线",\n'
            '  "primary_document": "最核心的文档名称",\n'
            '  "information_overlap": ["重叠的信息点1", "重叠的信息点2"]\n'
            "}"
        )
        try:
            return self.llm_client.call_llm(prompt, response_json=True)
        except Exception as e:
            print(f"    ⚠ 关联分析失败: {e}")
            return {"relationships": [], "overall_narrative": "", "recommended_storyline": ""}

    # ================================================================== #
    #  Layer 4 — 汇报框架映射
    # ================================================================== #
    def _map_to_framework(self, profiles: List[Dict], units: List[Dict],
                          relationships: Dict) -> tuple:
        """将语义单元映射到PPT汇报框架"""
        # 按颗粒度分组
        must_show = [u for u in units if u.get("granularity") == "must_show"]
        should_show = [u for u in units if u.get("granularity") == "should_show"]
        optional = [u for u in units if u.get("granularity") == "optional"]

        units_summary = {
            "must_show_count": len(must_show),
            "should_show_count": len(should_show),
            "optional_count": len(optional),
            "must_show_types": self._count_by_type(must_show),
            "should_show_types": self._count_by_type(should_show),
        }

        prompt = (
            "你是专业的PPT框架设计师。请根据以下信息，为PPT建立初步的框架映射。\n\n"
            f"文档画像: {json.dumps(profiles, ensure_ascii=False)}\n\n"
            f"语义单元统计: {json.dumps(units_summary, ensure_ascii=False)}\n\n"
            f"叙事主线: {relationships.get('recommended_storyline', '待定')}\n\n"
            "请输出JSON:\n"
            "{\n"
            '  "suggested_sections": [\n'
            '    {"title": "章节名", "purpose": "该章节的作用", "estimated_pages": 3, '
            '"content_types": ["achievement", "data"], "priority": "high/medium/low"}\n'
            "  ],\n"
            '  "content_actions": [\n'
            '    {"content": "信息描述", "action": "direct_use/transform/generate/skip", '
            '"reason": "原因"}\n'
            "  ],\n"
            '  "information_gaps": [\n'
            '    {"expected_info": "预期应有的信息", "status": "missing/partial/vague", '
            '"suggestion": "建议用户补充/Agent推断/跳过"}\n'
            "  ],\n"
            '  "estimated_total_pages": 15\n'
            "}"
        )
        try:
            result = self.llm_client.call_llm(prompt, response_json=True)
            gaps = result.get("information_gaps", [])
            return result, gaps
        except Exception as e:
            print(f"    ⚠ 框架映射失败: {e}")
            return {}, []

    # ================================================================== #
    #  辅助方法
    # ================================================================== #
    def _chunk_document(self, doc: Dict) -> List[Dict]:
        """按文档自然分节进行Chunk"""
        chunks = []
        content = doc["content"]
        file_type = doc["file_type"]

        if file_type == ".docx":
            # 按 Heading 分块
            current_chunk = {"section": "开头", "texts": []}
            for para in content.get("paragraphs", []):
                if para.get("is_heading") and para.get("heading_level", 0) <= 2:
                    if current_chunk["texts"]:
                        chunks.append({
                            "section": current_chunk["section"],
                            "text": "\n".join(current_chunk["texts"]),
                        })
                    current_chunk = {"section": para["text"], "texts": []}
                current_chunk["texts"].append(para["text"])
            if current_chunk["texts"]:
                chunks.append({
                    "section": current_chunk["section"],
                    "text": "\n".join(current_chunk["texts"]),
                })
            # 表格作为独立chunk
            for tbl in content.get("tables", []):
                table_text = "\n".join([" | ".join(row) for row in tbl["data"]])
                chunks.append({"section": f"表格{tbl['index']+1}", "text": table_text})

        elif file_type == ".pdf":
            # 按页分块，相邻页合并（避免过小）
            pages = content.get("pages", [])
            i = 0
            while i < len(pages):
                batch = pages[i:i+2]
                text = "\n".join(p["text"] for p in batch)
                chunks.append({"section": f"第{batch[0]['page_num']}页", "text": text})
                i += 2

        elif file_type == ".pptx":
            for slide in content.get("slides", []):
                parts = [slide.get("title", "")]
                parts.extend(slide.get("texts", []))
                if slide.get("notes"):
                    parts.append(f"[备注] {slide['notes']}")
                chunks.append({
                    "section": f"第{slide['slide_num']}页: {slide.get('title', '')}",
                    "text": "\n".join(parts),
                })

        else:
            # txt / md — 按段落分块
            text = content.get("text", "")
            paragraphs = text.split("\n\n")
            for i in range(0, len(paragraphs), 3):
                batch = paragraphs[i:i+3]
                chunks.append({"section": f"段落{i+1}", "text": "\n\n".join(batch)})

        # 如果没有分出chunk，整体作为一个chunk
        if not chunks:
            chunks.append({"section": "全文", "text": self._get_text_preview(doc, max_chars=6000)})

        return chunks

    def _get_text_preview(self, doc: Dict, max_chars: int = 3000) -> str:
        """获取文档的文本预览"""
        content = doc["content"]
        file_type = doc["file_type"]
        texts = []

        if file_type == ".docx":
            for para in content.get("paragraphs", []):
                texts.append(para["text"])
        elif file_type == ".pdf":
            for page in content.get("pages", []):
                texts.append(page["text"])
        elif file_type == ".pptx":
            for slide in content.get("slides", []):
                if slide.get("title"):
                    texts.append(slide["title"])
                texts.extend(slide.get("texts", []))
        else:
            texts.append(content.get("text", ""))

        full_text = "\n".join(texts)
        return full_text[:max_chars]

    def _extract_sections_from_paragraphs(self, paragraphs: List[Dict]) -> List[Dict]:
        """从段落中提取章节结构"""
        sections = []
        for para in paragraphs:
            if para.get("is_heading"):
                sections.append({
                    "title": para["text"],
                    "level": para.get("heading_level", 1),
                    "index": para["index"],
                })
        return sections

    def _merge_documents(self, documents: List[Dict]) -> str:
        """合并所有文档的文本内容"""
        parts = []
        for doc in documents:
            parts.append(f"\n{'='*40}\n文件: {doc['file_name']}\n{'='*40}")
            parts.append(self._get_text_preview(doc, max_chars=10000))
        return "\n".join(parts)

    def _count_by_type(self, units: List[Dict]) -> Dict[str, int]:
        """按类型统计语义单元数量"""
        counts = {}
        for u in units:
            t = u.get("type", "other")
            counts[t] = counts.get(t, 0) + 1
        return counts

    def _generate_summary(self, documents: List[Dict], semantic_units: List[Dict]) -> str:
        """生成解析摘要"""
        total_files = len(documents)
        type_counts = {}
        for doc in documents:
            ft = doc["file_type"]
            type_counts[ft] = type_counts.get(ft, 0) + 1

        parts = [f"共解析 {total_files} 个文件"]
        for ft, cnt in type_counts.items():
            parts.append(f"{cnt} 个 {ft} 文件")

        if semantic_units:
            must = sum(1 for u in semantic_units if u.get("granularity") == "must_show")
            parts.append(f"提取 {len(semantic_units)} 个语义单元（{must} 个必须呈现）")

        return "，".join(parts)

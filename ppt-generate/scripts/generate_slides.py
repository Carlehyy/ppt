#!/usr/bin/env python3
"""
Slide Generator Module — 逐页生成与修改请求处理
PRD要求：
1. 逐页生成 — 根据大纲和语义单元，为每页生成精确内容
2. 版式适配 — 内容严格适配版式约束
3. 溯源标注 — 每条信息标注来源和置信度
4. 修改请求处理 — 支持用户对特定页面的修改需求
"""

import os
import json
import copy
from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN


class SlideGenerator:
    """PPT幻灯片生成器"""

    def __init__(self, llm_client=None):
        self.llm_client = llm_client

    def generate(self, outline: Dict, parsed_content: Dict,
                 template_analysis: Dict, final_config: Dict,
                 template_path: str, output_path: str) -> Dict[str, Any]:
        """
        执行完整的PPT生成流程

        Returns:
            {
                "output_path": str,
                "total_slides": int,
                "slides_data": [],
                "generation_log": [],
            }
        """
        prs = Presentation(template_path)
        slides_data = []
        generation_log = []

        # 删除模板中的示例幻灯片
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].get(
                '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
            )
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

        sections = outline.get("sections", [])
        design_lang = template_analysis.get("design_language", {})
        content_guidelines = design_lang.get("content_guidelines", {})
        semantic_units = parsed_content.get("semantic_units", [])

        page_count = 0
        for section in sections:
            for page_spec in section.get("pages", []):
                page_count += 1
                page_num = page_spec.get("page_num", page_count)
                print(f"    生成第 {page_num} 页: {page_spec.get('title', '')}")

                try:
                    # Step 1: 用LLM生成页面内容
                    slide_content = self._generate_slide_content(
                        page_spec, semantic_units, content_guidelines, final_config
                    )

                    # Step 2: 选择版式并创建幻灯片
                    layout_index = page_spec.get("layout_index", 0)
                    layout_index = min(layout_index, len(prs.slide_layouts) - 1)
                    layout = prs.slide_layouts[layout_index]
                    slide = prs.slides.add_slide(layout)

                    # Step 3: 填充内容到幻灯片
                    self._fill_slide(slide, slide_content, page_spec)

                    # Step 4: 添加演讲备注
                    self._add_speaker_notes(slide, slide_content, page_spec)

                    slides_data.append({
                        "page_num": page_num,
                        "title": slide_content.get("title", ""),
                        "content_type": page_spec.get("content_type", ""),
                        "layout_index": layout_index,
                        "content": slide_content,
                        "status": "success",
                    })
                    generation_log.append({
                        "page": page_num,
                        "status": "success",
                        "message": f"成功生成: {slide_content.get('title', '')}",
                    })

                except Exception as e:
                    print(f"    ⚠ 第 {page_num} 页生成失败: {e}")
                    layout = prs.slide_layouts[0]
                    slide = prs.slides.add_slide(layout)
                    if slide.shapes.title:
                        slide.shapes.title.text = page_spec.get("title", f"第{page_num}页")

                    slides_data.append({
                        "page_num": page_num,
                        "title": page_spec.get("title", ""),
                        "content_type": page_spec.get("content_type", ""),
                        "status": "failed",
                        "error": str(e),
                    })
                    generation_log.append({
                        "page": page_num,
                        "status": "failed",
                        "message": str(e),
                    })

        # 保存PPT
        prs.save(output_path)
        print(f"  ✓ PPT已保存: {output_path}")

        return {
            "output_path": output_path,
            "total_slides": len(slides_data),
            "slides_data": slides_data,
            "generation_log": generation_log,
        }

    def _generate_slide_content(self, page_spec: Dict, semantic_units: List[Dict],
                                 content_guidelines: Dict, config: Dict) -> Dict:
        """使用LLM为单页生成内容"""
        content_type = page_spec.get("content_type", "")

        # 封面和结束页不需要LLM
        if content_type == "cover":
            return {
                "title": config.get("presentation_title", page_spec.get("title", "")),
                "subtitle": page_spec.get("subtitle", ""),
                "body": [],
                "notes": page_spec.get("speaker_notes", ""),
            }
        if content_type == "ending":
            return {
                "title": page_spec.get("title", "谢谢"),
                "subtitle": page_spec.get("subtitle", ""),
                "body": [],
                "notes": "",
            }
        if content_type == "section_divider":
            return {
                "title": page_spec.get("title", ""),
                "subtitle": page_spec.get("core_message", ""),
                "body": [],
                "notes": page_spec.get("speaker_notes", ""),
            }

        if not self.llm_client:
            return {
                "title": page_spec.get("title", ""),
                "body": page_spec.get("key_points", ["待填充"]),
                "notes": "",
                "source_info": [],
            }

        # 收集相关的语义单元
        relevant_units = self._find_relevant_units(page_spec, semantic_units)

        prompt = (
            "你是专业的PPT内容撰写专家。请为以下幻灯片生成精确的内容。\n\n"
            f"## 页面规格\n"
            f"- 标题: {page_spec.get('title', '')}\n"
            f"- 内容类型: {content_type}\n"
            f"- 核心信息: {page_spec.get('core_message', '')}\n"
            f"- 预期要点: {json.dumps(page_spec.get('key_points', []), ensure_ascii=False)}\n"
            f"- 相关数据: {json.dumps(page_spec.get('data_to_show', []), ensure_ascii=False)}\n\n"
            f"## 内容指南\n"
            f"- 标题风格: {content_guidelines.get('title_style', '简洁有力')}\n"
            f"- 要点风格: {content_guidelines.get('bullet_style', '短句')}\n"
            f"- 数据呈现: {content_guidelines.get('data_presentation', '数字突出')}\n"
            f"- 最大文字密度: {content_guidelines.get('max_text_density', '每页不超过5个要点')}\n"
            f"- 语言风格: {config.get('language_style', '专业')}\n\n"
            f"## 可用素材（语义单元）\n"
            f"{json.dumps(relevant_units, ensure_ascii=False)}\n\n"
            "## 生成规则\n"
            "1. 标题不超过15个字，简洁有力\n"
            "2. 正文要点3-5个，每个不超过30字\n"
            "3. 优先使用原始素材中的具体数据\n"
            "4. 每条信息标注来源和置信度\n"
            "5. 如果是对比类型，生成左右两列内容\n\n"
            "请输出JSON:\n"
            "{\n"
            '  "title": "页面标题",\n'
            '  "subtitle": "副标题（可选，留空则无）",\n'
            '  "body": ["要点1", "要点2", "要点3"],\n'
            '  "body_right": ["右列要点1", "右列要点2"],\n'
            '  "notes": "演讲备注",\n'
            '  "source_info": [\n'
            '    {"content": "要点1", "source": "文件:位置", "confidence": "high/medium/low"}\n'
            '  ],\n'
            '  "data_highlights": ["需要突出的数据"],\n'
            '  "visual_suggestion": "视觉建议（如建议插入图表）"\n'
            "}"
        )

        try:
            result = self.llm_client.call_llm(prompt, response_json=True)
            return result
        except Exception as e:
            print(f"      ⚠ LLM内容生成失败: {e}")
            return {
                "title": page_spec.get("title", ""),
                "body": page_spec.get("key_points", ["待填充"]),
                "notes": "",
                "source_info": [],
            }

    def _find_relevant_units(self, page_spec: Dict, semantic_units: List[Dict]) -> List[Dict]:
        """查找与当前页面相关的语义单元"""
        relevant = []
        content_type = page_spec.get("content_type", "")
        key_points = page_spec.get("key_points", [])

        type_mapping = {
            "achievement_list": ["achievement", "data"],
            "data_showcase": ["data", "achievement"],
            "comparison": ["data", "achievement", "problem"],
            "problem_analysis": ["problem", "method"],
            "plan_timeline": ["plan"],
            "conclusion": ["conclusion", "achievement"],
        }
        target_types = type_mapping.get(content_type, [])

        for unit in semantic_units:
            if unit.get("type") in target_types:
                relevant.append({
                    "content": unit.get("content", ""),
                    "type": unit.get("type"),
                    "source": unit.get("source", ""),
                    "confidence": unit.get("confidence", "medium"),
                    "key_data": unit.get("key_data", []),
                })
                continue

            content = unit.get("content", "")
            for kp in key_points:
                if len(kp) > 7 and kp[:8] in content:
                    relevant.append({
                        "content": content,
                        "type": unit.get("type"),
                        "source": unit.get("source", ""),
                        "confidence": unit.get("confidence", "medium"),
                        "key_data": unit.get("key_data", []),
                    })
                    break

        return relevant[:10]

    def _fill_slide(self, slide, content: Dict, page_spec: Dict):
        """将生成的内容填充到幻灯片"""
        content_type = page_spec.get("content_type", "")

        # 填充标题
        title_text = content.get("title", "")
        if slide.shapes.title and title_text:
            slide.shapes.title.text = title_text

        # 填充副标题
        subtitle_text = content.get("subtitle", "")
        if subtitle_text:
            for shape in slide.placeholders:
                ph_idx = shape.placeholder_format.idx
                if ph_idx == 1:  # 通常是副标题占位符
                    shape.text = subtitle_text
                    break

        # 填充正文内容
        body_items = content.get("body", [])
        if body_items:
            body_placeholder = self._find_body_placeholder(slide)
            if body_placeholder:
                self._fill_body_content(body_placeholder, body_items)

        # 对比类型：填充右列
        body_right = content.get("body_right", [])
        if body_right and content_type == "comparison":
            right_placeholder = self._find_right_placeholder(slide)
            if right_placeholder:
                self._fill_body_content(right_placeholder, body_right)

        # 处理表格数据
        table_data = content.get("table_data", [])
        if table_data:
            self._insert_table(slide, table_data)

    def _find_body_placeholder(self, slide):
        """查找正文占位符"""
        for shape in slide.placeholders:
            ph_idx = shape.placeholder_format.idx
            if ph_idx >= 1 and shape != slide.shapes.title:
                return shape
        return None

    def _find_right_placeholder(self, slide):
        """查找右侧占位符（用于对比版式）"""
        body_found = False
        for shape in slide.placeholders:
            ph_idx = shape.placeholder_format.idx
            if ph_idx >= 1 and shape != slide.shapes.title:
                if body_found:
                    return shape
                body_found = True
        return None

    def _fill_body_content(self, placeholder, items: List[str]):
        """填充正文内容到占位符"""
        tf = placeholder.text_frame
        tf.clear()

        for i, item in enumerate(items):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            para.text = str(item)
            para.level = 0

    def _insert_table(self, slide, table_data: List[List[str]]):
        """插入表格"""
        if not table_data:
            return
        rows = len(table_data)
        cols = len(table_data[0]) if rows > 0 else 0
        if rows == 0 or cols == 0:
            return

        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(0.5 * rows)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        for i, row_data in enumerate(table_data):
            for j, cell_value in enumerate(row_data):
                if j < cols:
                    table.cell(i, j).text = str(cell_value)

    def _add_speaker_notes(self, slide, content: Dict, page_spec: Dict):
        """添加演讲备注"""
        notes_parts = []

        notes_text = content.get("notes", "") or page_spec.get("speaker_notes", "")
        if notes_text:
            notes_parts.append(f"【演讲要点】\n{notes_text}")

        source_info = content.get("source_info", [])
        if source_info:
            notes_parts.append("\n【信息来源】")
            for si in source_info:
                conf_icon = {"high": "✅", "medium": "⚡", "low": "❓"}.get(
                    si.get("confidence", "medium"), "⚡"
                )
                notes_parts.append(
                    f"  {conf_icon} {si.get('content', '')[:30]} — 来源: {si.get('source', '未知')}"
                )

        visual = content.get("visual_suggestion", "")
        if visual:
            notes_parts.append(f"\n【视觉建议】\n{visual}")

        if notes_parts:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = "\n".join(notes_parts)

    # ================================================================== #
    #  修改请求处理
    # ================================================================== #
    def handle_modification_request(self, modification: Dict, slides_data: List[Dict],
                                     outline: Dict, parsed_content: Dict,
                                     template_analysis: Dict, config: Dict,
                                     template_path: str, output_path: str) -> Dict:
        """
        处理用户的修改请求

        Args:
            modification: {
                "type": "modify_page/add_page/delete_page/reorder/global_style",
                "target_page": 5,
                "instruction": "用户的修改指令",
            }
        """
        mod_type = modification.get("type", "")
        print(f"  [修改] 处理修改请求: {mod_type}")

        if mod_type == "modify_page":
            return self._modify_single_page(
                modification, slides_data, parsed_content,
                template_analysis, config, template_path, output_path
            )
        elif mod_type == "add_page":
            return self._add_page(
                modification, slides_data, outline, parsed_content,
                template_analysis, config, template_path, output_path
            )
        elif mod_type == "delete_page":
            return self._delete_page(
                modification, slides_data, template_path, output_path
            )
        elif mod_type == "global_style":
            return self._apply_global_style_change(
                modification, slides_data, parsed_content,
                template_analysis, config, template_path, output_path
            )
        else:
            return {"status": "error", "message": f"不支持的修改类型: {mod_type}"}

    def _modify_single_page(self, modification, slides_data, parsed_content,
                             template_analysis, config, template_path, output_path):
        """修改单页内容"""
        target_page = modification.get("target_page", 1)
        instruction = modification.get("instruction", "")

        target_data = None
        for sd in slides_data:
            if sd.get("page_num") == target_page:
                target_data = sd
                break

        if not target_data:
            return {"status": "error", "message": f"未找到第{target_page}页"}

        if self.llm_client:
            prompt = (
                "你是PPT内容修改专家。请根据用户指令修改以下页面内容。\n\n"
                f"## 当前内容\n{json.dumps(target_data.get('content', {}), ensure_ascii=False)}\n\n"
                f"## 用户修改指令\n{instruction}\n\n"
                "请输出修改后的完整JSON（格式与当前内容相同）。"
            )
            try:
                new_content = self.llm_client.call_llm(prompt, response_json=True)
                target_data["content"] = new_content
                target_data["status"] = "modified"
            except Exception as e:
                return {"status": "error", "message": f"修改失败: {e}"}

        return self._regenerate_pptx(slides_data, template_analysis, config, template_path, output_path)

    def _add_page(self, modification, slides_data, outline, parsed_content,
                   template_analysis, config, template_path, output_path):
        """添加新页面"""
        after_page = modification.get("target_page", len(slides_data))
        instruction = modification.get("instruction", "")

        new_page_spec = {
            "page_num": after_page + 1,
            "title": instruction[:15] if instruction else "新页面",
            "content_type": "achievement_list",
            "layout_index": 1,
            "core_message": instruction,
            "key_points": [],
        }

        if self.llm_client:
            content = self._generate_slide_content(
                new_page_spec, parsed_content.get("semantic_units", []),
                template_analysis.get("design_language", {}).get("content_guidelines", {}),
                config
            )
        else:
            content = {"title": new_page_spec["title"], "body": ["待填充"]}

        new_slide_data = {
            "page_num": after_page + 1,
            "title": content.get("title", ""),
            "content_type": new_page_spec["content_type"],
            "layout_index": new_page_spec["layout_index"],
            "content": content,
            "status": "added",
        }

        insert_idx = after_page
        slides_data.insert(insert_idx, new_slide_data)
        for i, sd in enumerate(slides_data):
            sd["page_num"] = i + 1

        return self._regenerate_pptx(slides_data, template_analysis, config, template_path, output_path)

    def _delete_page(self, modification, slides_data, template_path, output_path):
        """删除页面"""
        target_page = modification.get("target_page", 0)
        slides_data[:] = [sd for sd in slides_data if sd.get("page_num") != target_page]
        for i, sd in enumerate(slides_data):
            sd["page_num"] = i + 1
        return {"status": "success", "slides_data": slides_data,
                "message": f"已删除第{target_page}页"}

    def _apply_global_style_change(self, modification, slides_data, parsed_content,
                                    template_analysis, config, template_path, output_path):
        """应用全局风格修改"""
        instruction = modification.get("instruction", "")

        if self.llm_client:
            for sd in slides_data:
                if sd.get("content_type") in ("cover", "ending", "section_divider"):
                    continue
                content = sd.get("content", {})
                prompt = (
                    f"请根据以下风格指令调整内容:\n"
                    f"指令: {instruction}\n"
                    f"当前内容: {json.dumps(content, ensure_ascii=False)}\n"
                    "输出调整后的JSON。"
                )
                try:
                    new_content = self.llm_client.call_llm(prompt, response_json=True)
                    sd["content"] = new_content
                    sd["status"] = "style_modified"
                except Exception:
                    pass

        return self._regenerate_pptx(slides_data, template_analysis, config, template_path, output_path)

    def _regenerate_pptx(self, slides_data, template_analysis, config,
                          template_path, output_path):
        """根据slides_data重新生成PPTX文件"""
        prs = Presentation(template_path)

        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].get(
                '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
            )
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

        for sd in slides_data:
            layout_index = sd.get("layout_index", 0)
            layout_index = min(layout_index, len(prs.slide_layouts) - 1)
            layout = prs.slide_layouts[layout_index]
            slide = prs.slides.add_slide(layout)

            content = sd.get("content", {})
            page_spec = {"content_type": sd.get("content_type", "")}
            self._fill_slide(slide, content, page_spec)
            self._add_speaker_notes(slide, content, page_spec)

        prs.save(output_path)
        return {
            "status": "success",
            "output_path": output_path,
            "total_slides": len(slides_data),
            "slides_data": slides_data,
        }

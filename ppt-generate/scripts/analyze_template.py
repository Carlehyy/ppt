#!/usr/bin/env python3
"""
Template Analyzer Module — 四层模板分析架构
Layer 1: 视觉层 — 配色方案、字体体系、背景风格、图形元素
Layer 2: 版式层 — 每种类型页面的布局，文字/图表/留白比例关系
Layer 3: 结构层 — 模板暗含的叙事节奏
Layer 4: 设计语言层 — 模板传递的"气质"（LLM驱动）
"""

import os
import json
from typing import Dict, Any, List, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree


class TemplateAnalyzer:
    """四层模板分析器"""

    def __init__(self, llm_client=None):
        self.llm_client = llm_client

    def analyze(self, template_path: str) -> Dict[str, Any]:
        """
        执行完整的四层模板分析

        Returns:
            {
                "template_path": str,
                "template_name": str,
                "visual_profile": {},    # Layer 1 — 视觉层
                "layouts": [],           # Layer 2 — 版式层
                "structure_profile": {}, # Layer 3 — 结构层
                "design_language": {},   # Layer 4 — 设计语言层
                "total_layouts": int,
                "layout_recommendations": {}  # 版式使用建议
            }
        """
        print(f"  分析模板: {os.path.basename(template_path)}")
        prs = Presentation(template_path)

        # ---- Layer 1: 视觉层 ---- #
        print("  [Layer 1] 视觉层分析...")
        visual_profile = self._analyze_visual_layer(prs)
        print("  [Layer 1] ✓ 完成")

        # ---- Layer 2: 版式层 ---- #
        print("  [Layer 2] 版式层分析...")
        layouts = self._analyze_layout_layer(prs)
        print(f"  [Layer 2] ✓ 完成，发现 {len(layouts)} 个版式")

        # ---- Layer 3: 结构层 ---- #
        print("  [Layer 3] 结构层分析...")
        structure_profile = self._analyze_structure_layer(layouts)
        print("  [Layer 3] ✓ 完成")

        # ---- Layer 4: 设计语言层（LLM驱动） ---- #
        design_language = {}
        layout_recommendations = {}
        if self.llm_client:
            print("  [Layer 4] 设计语言层分析...")
            design_language = self._analyze_design_language(visual_profile, layouts, structure_profile)
            layout_recommendations = self._generate_layout_recommendations(layouts, design_language)
            print("  [Layer 4] ✓ 完成")
        else:
            print("  [Layer 4] ⚠ 未配置LLM，跳过设计语言分析")

        return {
            "template_path": template_path,
            "template_name": os.path.basename(template_path),
            "visual_profile": visual_profile,
            "layouts": layouts,
            "structure_profile": structure_profile,
            "design_language": design_language,
            "total_layouts": len(layouts),
            "layout_recommendations": layout_recommendations,
        }

    # ================================================================== #
    #  Layer 1 — 视觉层分析
    # ================================================================== #
    def _analyze_visual_layer(self, prs: Presentation) -> Dict[str, Any]:
        """提取配色方案、字体体系、背景风格"""
        profile = {
            "slide_width": prs.slide_width,
            "slide_height": prs.slide_height,
            "slide_width_inches": round(prs.slide_width / 914400, 2),
            "slide_height_inches": round(prs.slide_height / 914400, 2),
            "color_scheme": self._extract_color_scheme(prs),
            "font_scheme": self._extract_font_scheme(prs),
            "background_style": self._extract_background_style(prs),
            "colors_used": self._collect_colors_from_slides(prs),
            "fonts_used": self._collect_fonts_from_slides(prs),
        }
        return profile

    def _extract_color_scheme(self, prs: Presentation) -> Dict[str, Any]:
        """从主题中提取配色方案"""
        colors = {}
        try:
            theme = prs.slide_masters[0].element.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}theme"
            )
            if theme is None:
                # 尝试从slide_master的XML中提取
                master_xml = prs.slide_masters[0].element
                ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
                color_elements = master_xml.findall(".//a:clrScheme/*", ns)
                for elem in color_elements:
                    tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
                    # 尝试获取颜色值
                    for child in elem:
                        if "val" in child.attrib:
                            colors[tag] = child.attrib["val"]
                        elif "lastClr" in child.attrib:
                            colors[tag] = child.attrib["lastClr"]
        except Exception:
            pass

        # 如果主题提取失败，尝试从实际幻灯片中收集
        if not colors:
            colors = {"note": "主题颜色提取受限，已从实际内容中采集颜色"}

        return colors

    def _extract_font_scheme(self, prs: Presentation) -> Dict[str, Any]:
        """从主题中提取字体体系"""
        fonts = {"major": "", "minor": ""}
        try:
            ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
            master_xml = prs.slide_masters[0].element
            major_font = master_xml.find(".//a:majorFont/a:latin", ns)
            minor_font = master_xml.find(".//a:minorFont/a:latin", ns)
            if major_font is not None:
                fonts["major"] = major_font.attrib.get("typeface", "")
            if minor_font is not None:
                fonts["minor"] = minor_font.attrib.get("typeface", "")

            # 尝试获取东亚字体
            major_ea = master_xml.find(".//a:majorFont/a:ea", ns)
            minor_ea = master_xml.find(".//a:minorFont/a:ea", ns)
            if major_ea is not None:
                fonts["major_ea"] = major_ea.attrib.get("typeface", "")
            if minor_ea is not None:
                fonts["minor_ea"] = minor_ea.attrib.get("typeface", "")
        except Exception:
            pass
        return fonts

    def _extract_background_style(self, prs: Presentation) -> Dict[str, Any]:
        """分析背景风格"""
        bg_info = {"type": "unknown", "description": ""}
        try:
            if prs.slide_masters:
                master = prs.slide_masters[0]
                bg = master.background
                if bg and bg.fill:
                    fill = bg.fill
                    if fill.type is not None:
                        bg_info["type"] = str(fill.type)
        except Exception:
            pass
        return bg_info

    def _collect_colors_from_slides(self, prs: Presentation) -> List[str]:
        """从实际幻灯片内容中收集使用的颜色"""
        colors_found = set()
        try:
            for layout in prs.slide_layouts:
                for shape in layout.placeholders:
                    if hasattr(shape, "text_frame"):
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.color and run.font.color.rgb:
                                    colors_found.add(str(run.font.color.rgb))
            # 也从已有幻灯片中收集
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                try:
                                    if run.font.color and run.font.color.rgb:
                                        colors_found.add(str(run.font.color.rgb))
                                except Exception:
                                    pass
        except Exception:
            pass
        return list(colors_found)[:20]

    def _collect_fonts_from_slides(self, prs: Presentation) -> List[str]:
        """从实际内容中收集使用的字体"""
        fonts_found = set()
        try:
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.name:
                                    fonts_found.add(run.font.name)
            for layout in prs.slide_layouts:
                for shape in layout.placeholders:
                    if hasattr(shape, "text_frame"):
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.name:
                                    fonts_found.add(run.font.name)
        except Exception:
            pass
        return list(fonts_found)[:15]

    # ================================================================== #
    #  Layer 2 — 版式层分析
    # ================================================================== #
    def _analyze_layout_layer(self, prs: Presentation) -> List[Dict[str, Any]]:
        """分析所有版式的详细布局信息"""
        layouts = []
        for idx, layout in enumerate(prs.slide_layouts):
            placeholders = self._analyze_placeholders(layout)

            # 计算布局比例
            content_area = self._calculate_content_area(placeholders, prs)

            layout_info = {
                "index": idx,
                "name": layout.name,
                "placeholders": placeholders,
                "total_placeholders": len(placeholders),
                "content_area": content_area,
                "layout_category": self._categorize_layout(layout.name, placeholders),
                "best_for": self._infer_layout_usage(layout.name, placeholders),
                "capacity": self._estimate_content_capacity(placeholders),
            }
            layouts.append(layout_info)
        return layouts

    def _analyze_placeholders(self, layout) -> List[Dict[str, Any]]:
        """详细分析版式中的占位符"""
        placeholders = []
        type_map = {
            0: "TITLE", 1: "BODY", 2: "CENTER_TITLE", 3: "SUBTITLE",
            4: "DATE", 5: "SLIDE_NUMBER", 6: "FOOTER", 7: "HEADER",
            12: "PICTURE", 13: "CHART", 14: "TABLE", 15: "CLIP_ART",
            16: "MEDIA_CLIP", 17: "OBJECT", 18: "VERTICAL_BODY",
        }

        for ph in layout.placeholders:
            ph_type = ph.placeholder_format.type
            # python-pptx 的 type 可能是 int 或 enum
            type_val = ph_type if isinstance(ph_type, int) else (ph_type.real if ph_type else -1)

            ph_info = {
                "idx": ph.placeholder_format.idx,
                "type": type_val,
                "type_name": type_map.get(type_val, f"TYPE_{type_val}"),
                "name": ph.name,
                "left": ph.left,
                "top": ph.top,
                "width": ph.width,
                "height": ph.height,
                "left_inches": round(ph.left / 914400, 2) if ph.left else 0,
                "top_inches": round(ph.top / 914400, 2) if ph.top else 0,
                "width_inches": round(ph.width / 914400, 2) if ph.width else 0,
                "height_inches": round(ph.height / 914400, 2) if ph.height else 0,
            }

            # 分析文本框的默认样式
            if hasattr(ph, "text_frame"):
                try:
                    tf = ph.text_frame
                    if tf.paragraphs:
                        first_para = tf.paragraphs[0]
                        ph_info["default_alignment"] = str(first_para.alignment) if first_para.alignment else "LEFT"
                        for run in first_para.runs:
                            if run.font.size:
                                ph_info["default_font_size_pt"] = round(run.font.size / 12700, 1)
                            if run.font.bold:
                                ph_info["default_bold"] = True
                            break
                except Exception:
                    pass

            placeholders.append(ph_info)
        return placeholders

    def _calculate_content_area(self, placeholders: List[Dict], prs: Presentation) -> Dict:
        """计算内容区域占比"""
        slide_area = (prs.slide_width / 914400) * (prs.slide_height / 914400)
        content_phs = [p for p in placeholders if p["type_name"] in ("TITLE", "BODY", "CENTER_TITLE", "SUBTITLE")]
        content_area = sum(p["width_inches"] * p["height_inches"] for p in content_phs)

        return {
            "slide_area_sq_inches": round(slide_area, 2),
            "content_area_sq_inches": round(content_area, 2),
            "content_ratio": round(content_area / slide_area, 2) if slide_area > 0 else 0,
            "whitespace_ratio": round(1 - content_area / slide_area, 2) if slide_area > 0 else 1,
        }

    def _categorize_layout(self, name: str, placeholders: List[Dict]) -> str:
        """将版式归类"""
        name_lower = name.lower()
        ph_types = {p["type_name"] for p in placeholders}

        if "title" in name_lower and ("slide" in name_lower or "CENTER_TITLE" in ph_types):
            return "cover"
        elif "section" in name_lower or "divider" in name_lower:
            return "section_divider"
        elif "two" in name_lower or "comparison" in name_lower:
            return "two_column"
        elif "picture" in name_lower or "PICTURE" in ph_types:
            return "image_content"
        elif "chart" in name_lower or "CHART" in ph_types:
            return "chart_content"
        elif "table" in name_lower or "TABLE" in ph_types:
            return "table_content"
        elif "blank" in name_lower:
            return "blank"
        elif "TITLE" in ph_types and "BODY" in ph_types:
            return "title_and_content"
        elif "TITLE" in ph_types:
            return "title_only"
        else:
            return "custom"

    def _infer_layout_usage(self, name: str, placeholders: List[Dict]) -> str:
        """推断版式的最佳使用场景"""
        category = self._categorize_layout(name, placeholders)
        usage_map = {
            "cover": "封面页、标题页",
            "section_divider": "章节分隔页、过渡页",
            "two_column": "对比分析、并列展示、优劣势对比",
            "image_content": "图片展示、视觉内容、案例展示",
            "chart_content": "数据图表、统计展示、趋势分析",
            "table_content": "数据表格、详细对比、清单列表",
            "blank": "自定义内容、灵活布局、全图页面",
            "title_and_content": "标准内容页、要点列举、成果展示、计划说明",
            "title_only": "标题页、过渡页、强调页",
            "custom": "自定义布局",
        }
        return usage_map.get(category, "通用内容页")

    def _estimate_content_capacity(self, placeholders: List[Dict]) -> Dict:
        """估算版式的内容容量"""
        body_phs = [p for p in placeholders if p["type_name"] in ("BODY", "SUBTITLE")]
        if not body_phs:
            return {"max_bullets": 0, "supports_table": False, "supports_image": False}

        # 根据body区域高度估算最大要点数
        max_height = max(p["height_inches"] for p in body_phs) if body_phs else 0
        max_bullets = max(1, int(max_height / 0.6))  # 每个要点约0.6英寸

        has_picture_ph = any(p["type_name"] == "PICTURE" for p in placeholders)
        has_table_ph = any(p["type_name"] == "TABLE" for p in placeholders)

        return {
            "max_bullets": min(max_bullets, 8),
            "supports_table": has_table_ph or max_height > 3,
            "supports_image": has_picture_ph,
        }

    # ================================================================== #
    #  Layer 3 — 结构层分析
    # ================================================================== #
    def _analyze_structure_layer(self, layouts: List[Dict]) -> Dict:
        """分析模板暗含的叙事节奏"""
        categories = [l["layout_category"] for l in layouts]

        # 统计版式类型分布
        category_counts = {}
        for c in categories:
            category_counts[c] = category_counts.get(c, 0) + 1

        # 判断模板的叙事节奏特征
        has_cover = "cover" in categories
        has_divider = "section_divider" in categories
        content_count = categories.count("title_and_content")
        comparison_count = categories.count("two_column")
        data_count = categories.count("chart_content") + categories.count("table_content")

        # 推断叙事节奏
        rhythm = "standard"
        if data_count >= 2:
            rhythm = "data_driven"
        elif comparison_count >= 2:
            rhythm = "comparison_heavy"
        elif content_count >= 5:
            rhythm = "content_rich"

        return {
            "category_distribution": category_counts,
            "has_cover": has_cover,
            "has_section_dividers": has_divider,
            "narrative_rhythm": rhythm,
            "rhythm_description": {
                "standard": "标准叙事节奏，适合大多数汇报场景",
                "data_driven": "数据驱动型，适合数据分析和指标汇报",
                "comparison_heavy": "对比分析型，适合方案比选和竞品分析",
                "content_rich": "内容丰富型，适合详细的工作汇报和项目总结",
            }.get(rhythm, "标准叙事节奏"),
            "suggested_section_pattern": self._suggest_section_pattern(has_cover, has_divider, categories),
        }

    def _suggest_section_pattern(self, has_cover: bool, has_divider: bool,
                                  categories: List[str]) -> List[str]:
        """建议的章节模式"""
        pattern = []
        if has_cover:
            pattern.append("cover")
        pattern.append("content")
        if has_divider:
            pattern.extend(["section_divider", "content"])
        pattern.append("ending")
        return pattern

    # ================================================================== #
    #  Layer 4 — 设计语言层（LLM驱动）
    # ================================================================== #
    def _analyze_design_language(self, visual: Dict, layouts: List[Dict],
                                  structure: Dict) -> Dict:
        """使用LLM分析模板的设计语言和气质"""
        # 准备分析上下文
        context = {
            "color_scheme": visual.get("color_scheme", {}),
            "colors_used": visual.get("colors_used", [])[:10],
            "font_scheme": visual.get("font_scheme", {}),
            "fonts_used": visual.get("fonts_used", [])[:10],
            "background_style": visual.get("background_style", {}),
            "layout_categories": structure.get("category_distribution", {}),
            "narrative_rhythm": structure.get("narrative_rhythm", ""),
            "total_layouts": len(layouts),
            "layout_names": [l["name"] for l in layouts],
        }

        prompt = (
            "你是专业的PPT设计顾问。请根据以下模板的视觉和结构信息，分析其设计语言和气质。\n\n"
            f"模板信息:\n{json.dumps(context, ensure_ascii=False)}\n\n"
            "请输出JSON:\n"
            "{\n"
            '  "design_temperament": "简约商务/科技现代/学术严谨/创意活泼/稳重大气/清新淡雅",\n'
            '  "formality_level": "非常正式/正式/半正式/轻松",\n'
            '  "visual_density": "密集/适中/留白较多",\n'
            '  "color_mood": "对配色传递的情绪的描述",\n'
            '  "typography_style": "对字体风格的描述",\n'
            '  "recommended_tone": "建议的文案语气描述",\n'
            '  "content_guidelines": {\n'
            '    "title_style": "标题应该是什么风格（如：简洁有力/描述性/问句式）",\n'
            '    "bullet_style": "要点应该是什么风格（如：短句/动词开头/数据驱动）",\n'
            '    "data_presentation": "数据呈现建议（如：图表优先/数字突出/对比展示）",\n'
            '    "max_text_density": "建议的最大文字密度（如：每页不超过5个要点）"\n'
            "  },\n"
            '  "do_list": ["应该做的设计决策1", "应该做的设计决策2"],\n'
            '  "dont_list": ["不应该做的设计决策1", "不应该做的设计决策2"]\n'
            "}"
        )
        try:
            return self.llm_client.call_llm(prompt, response_json=True)
        except Exception as e:
            print(f"    ⚠ 设计语言分析失败: {e}")
            return {"design_temperament": "未知", "formality_level": "正式"}

    def _generate_layout_recommendations(self, layouts: List[Dict],
                                          design_language: Dict) -> Dict:
        """为每种内容类型生成版式推荐"""
        layout_summary = []
        for l in layouts:
            layout_summary.append({
                "index": l["index"],
                "name": l["name"],
                "category": l["layout_category"],
                "best_for": l["best_for"],
                "capacity": l["capacity"],
            })

        prompt = (
            "你是专业的PPT版式顾问。请根据以下可用版式和设计语言，为不同的内容类型推荐最佳版式。\n\n"
            f"可用版式:\n{json.dumps(layout_summary, ensure_ascii=False)}\n\n"
            f"设计语言: {design_language.get('design_temperament', '未知')}\n\n"
            "请为以下内容类型推荐版式，输出JSON:\n"
            "{\n"
            '  "cover": {"layout_index": 0, "layout_name": "...", "reason": "..."},\n'
            '  "section_divider": {"layout_index": 2, "layout_name": "...", "reason": "..."},\n'
            '  "achievement_list": {"layout_index": 1, "layout_name": "...", "reason": "..."},\n'
            '  "data_showcase": {"layout_index": 1, "layout_name": "...", "reason": "..."},\n'
            '  "comparison": {"layout_index": 3, "layout_name": "...", "reason": "..."},\n'
            '  "problem_analysis": {"layout_index": 1, "layout_name": "...", "reason": "..."},\n'
            '  "plan_timeline": {"layout_index": 1, "layout_name": "...", "reason": "..."},\n'
            '  "conclusion": {"layout_index": 1, "layout_name": "...", "reason": "..."},\n'
            '  "ending": {"layout_index": 0, "layout_name": "...", "reason": "..."}\n'
            "}\n"
            "注意: layout_index 必须在 0 到 " + str(len(layouts) - 1) + " 之间"
        )
        try:
            return self.llm_client.call_llm(prompt, response_json=True)
        except Exception as e:
            print(f"    ⚠ 版式推荐生成失败: {e}")
            return {}

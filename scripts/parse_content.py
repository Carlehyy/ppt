#!/usr/bin/env python3
"""
Content Parser Module
Extracts and analyzes content from input documents (Docx, PDF, PPT)
"""

import os
from typing import List, Dict, Any
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
from .utils import get_file_extension, format_source_info


class ContentParser:
    """Parse content from various document formats"""
    
    def __init__(self, llm_client=None):
        self.llm_client = llm_client
    
    def parse(self, file_paths: List[str]) -> Dict[str, Any]:
        """
        Parse multiple input files and extract structured content
        
        Args:
            file_paths: List of file paths to parse
            
        Returns:
            Dictionary containing parsed content and semantic analysis
        """
        all_documents = []
        
        for file_path in file_paths:
            print(f"  解析文件: {os.path.basename(file_path)}")
            
            ext = get_file_extension(file_path)
            
            if ext == '.docx':
                content = self._parse_docx(file_path)
            elif ext == '.pdf':
                content = self._parse_pdf(file_path)
            elif ext == '.pptx':
                content = self._parse_pptx(file_path)
            else:
                print(f"  警告: 不支持的文件格式 {ext}")
                continue
            
            all_documents.append({
                'file_path': file_path,
                'file_name': os.path.basename(file_path),
                'file_type': ext,
                'content': content
            })
        
        # Merge and analyze content
        merged_content = self._merge_documents(all_documents)
        
        # Use LLM to extract semantic units if available
        if self.llm_client:
            semantic_units = self.llm_client.extract_semantic_units(merged_content)
        else:
            semantic_units = []
        
        return {
            'documents': all_documents,
            'merged_content': merged_content,
            'semantic_units': semantic_units,
            'summary': self._generate_summary(all_documents)
        }
    
    def _parse_docx(self, file_path: str) -> Dict[str, Any]:
        """Parse Word document"""
        doc = Document(file_path)
        
        paragraphs = []
        for i, para in enumerate(doc.paragraphs):
            if para.text.strip():
                paragraphs.append({
                    'index': i,
                    'text': para.text,
                    'style': para.style.name
                })
        
        tables = []
        for i, table in enumerate(doc.tables):
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            tables.append({
                'index': i,
                'data': table_data
            })
        
        return {
            'paragraphs': paragraphs,
            'tables': tables,
            'total_paragraphs': len(paragraphs),
            'total_tables': len(tables)
        }
    
    def _parse_pdf(self, file_path: str) -> Dict[str, Any]:
        """Parse PDF document"""
        reader = PdfReader(file_path)
        
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text.strip():
                pages.append({
                    'page_num': i + 1,
                    'text': text
                })
        
        return {
            'pages': pages,
            'total_pages': len(pages)
        }
    
    def _parse_pptx(self, file_path: str) -> Dict[str, Any]:
        """Parse PowerPoint presentation"""
        prs = Presentation(file_path)
        
        slides = []
        for i, slide in enumerate(prs.slides):
            slide_content = {
                'slide_num': i + 1,
                'title': '',
                'text': []
            }
            
            # Extract title
            if slide.shapes.title:
                slide_content['title'] = slide.shapes.title.text
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    if shape != slide.shapes.title:
                        slide_content['text'].append(shape.text)
            
            slides.append(slide_content)
        
        return {
            'slides': slides,
            'total_slides': len(slides)
        }
    
    def _merge_documents(self, documents: List[Dict]) -> str:
        """Merge content from all documents into a single text"""
        merged_text = []
        
        for doc in documents:
            merged_text.append(f"\n=== {doc['file_name']} ===\n")
            
            content = doc['content']
            
            if doc['file_type'] == '.docx':
                for para in content['paragraphs']:
                    merged_text.append(para['text'])
            
            elif doc['file_type'] == '.pdf':
                for page in content['pages']:
                    merged_text.append(page['text'])
            
            elif doc['file_type'] == '.pptx':
                for slide in content['slides']:
                    if slide['title']:
                        merged_text.append(f"## {slide['title']}")
                    merged_text.extend(slide['text'])
        
        return '\n'.join(merged_text)
    
    def _generate_summary(self, documents: List[Dict]) -> str:
        """Generate a summary of parsed documents"""
        total_files = len(documents)
        file_types = {}
        
        for doc in documents:
            file_type = doc['file_type']
            file_types[file_type] = file_types.get(file_type, 0) + 1
        
        summary_parts = [f"共解析 {total_files} 个文件"]
        for file_type, count in file_types.items():
            summary_parts.append(f"{count} 个 {file_type} 文件")
        
        return ", ".join(summary_parts)

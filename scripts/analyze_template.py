#!/usr/bin/env python3
"""
Template Analyzer Module
Analyzes PowerPoint templates to extract design DNA
"""

import os
from typing import Dict, Any, List
from pptx import Presentation
from pptx.util import Inches, Pt


class TemplateAnalyzer:
    """Analyze PowerPoint template structure and design"""
    
    def __init__(self, llm_client=None):
        self.llm_client = llm_client
    
    def analyze(self, template_path: str) -> Dict[str, Any]:
        """
        Analyze a PowerPoint template
        
        Args:
            template_path: Path to the template file
            
        Returns:
            Dictionary containing template profile
        """
        print(f"  分析模板: {os.path.basename(template_path)}")
        
        prs = Presentation(template_path)
        
        # Extract visual profile
        visual_profile = self._extract_visual_profile(prs)
        
        # Analyze all layouts
        layouts = self._analyze_layouts(prs)
        
        return {
            'template_path': template_path,
            'template_name': os.path.basename(template_path),
            'visual_profile': visual_profile,
            'layouts': layouts,
            'total_layouts': len(layouts)
        }
    
    def _extract_visual_profile(self, prs: Presentation) -> Dict[str, Any]:
        """Extract visual design elements from template"""
        visual_profile = {
            'slide_width': prs.slide_width,
            'slide_height': prs.slide_height,
            'theme_colors': {},
            'fonts': {}
        }
        
        # Try to extract theme colors
        try:
            if hasattr(prs, 'core_properties'):
                # Extract basic properties
                pass
        except:
            pass
        
        # Note: Extracting exact theme colors from python-pptx is complex
        # We'll rely on LLM analysis of layout screenshots for this
        
        return visual_profile
    
    def _analyze_layouts(self, prs: Presentation) -> List[Dict[str, Any]]:
        """Analyze all slide layouts in the template"""
        layouts = []
        
        for idx, layout in enumerate(prs.slide_layouts):
            layout_info = {
                'index': idx,
                'name': layout.name,
                'width': layout.width,
                'height': layout.height,
                'placeholders': self._analyze_placeholders(layout),
                'total_placeholders': len(layout.placeholders),
                'best_for': self._infer_layout_usage(layout.name, layout.placeholders)
            }
            
            layouts.append(layout_info)
        
        return layouts
    
    def _analyze_placeholders(self, layout) -> List[Dict[str, Any]]:
        """Analyze placeholders in a layout"""
        placeholders = []
        
        for placeholder in layout.placeholders:
            ph_info = {
                'idx': placeholder.placeholder_format.idx,
                'type': placeholder.placeholder_format.type,
                'name': placeholder.name,
                'left': placeholder.left,
                'top': placeholder.top,
                'width': placeholder.width,
                'height': placeholder.height
            }
            
            # Map type number to readable name
            type_map = {
                1: 'TITLE',
                2: 'BODY',
                3: 'CENTER_TITLE',
                4: 'SUBTITLE',
                5: 'DATE',
                6: 'SLIDE_NUMBER',
                7: 'FOOTER',
                8: 'HEADER',
                13: 'PICTURE',
                14: 'CHART',
                15: 'TABLE',
                16: 'CLIP_ART',
                18: 'OBJECT'
            }
            
            ph_info['type_name'] = type_map.get(placeholder.placeholder_format.type, 'UNKNOWN')
            
            placeholders.append(ph_info)
        
        return placeholders
    
    def _infer_layout_usage(self, layout_name: str, placeholders: List[Dict]) -> str:
        """Infer best usage scenario for a layout based on name and structure"""
        name_lower = layout_name.lower()
        
        # Common layout patterns
        if 'title' in name_lower and 'slide' in name_lower:
            return "封面页、章节标题页"
        elif 'title' in name_lower and 'content' in name_lower:
            return "标准内容页、要点列举、成果展示"
        elif 'section' in name_lower or 'divider' in name_lower:
            return "章节分隔页"
        elif 'two' in name_lower or 'comparison' in name_lower:
            return "对比分析、并列展示"
        elif 'blank' in name_lower:
            return "自定义内容、灵活布局"
        elif 'picture' in name_lower or 'image' in name_lower:
            return "图片展示、视觉内容"
        elif 'chart' in name_lower or 'data' in name_lower:
            return "数据图表、统计展示"
        elif 'quote' in name_lower:
            return "引用、重点强调"
        elif 'thank' in name_lower or 'end' in name_lower:
            return "结束页、致谢页"
        else:
            # Infer from placeholder structure
            has_title = any(ph['type_name'] == 'TITLE' for ph in placeholders)
            has_body = any(ph['type_name'] == 'BODY' for ph in placeholders)
            
            if has_title and has_body:
                return "通用内容页"
            elif has_title:
                return "标题页、过渡页"
            else:
                return "自定义布局"

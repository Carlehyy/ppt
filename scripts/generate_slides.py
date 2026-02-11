#!/usr/bin/env python3
"""
Slides Generator Module
Generates PowerPoint presentations from structured content
"""

import os
from typing import Dict, List, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


class SlidesGenerator:
    """Generate PowerPoint slides from structured data"""
    
    def __init__(self, config: Dict = None):
        self.config = config or {}
    
    def generate(self, outline: Dict, template_path: str, slides_content: List[Dict]) -> str:
        """
        Generate PowerPoint presentation
        
        Args:
            outline: Presentation outline structure
            template_path: Path to template file
            slides_content: List of slide content dictionaries
            
        Returns:
            Path to generated PPTX file
        """
        print("  生成PowerPoint文件...")
        
        # Load template
        prs = Presentation(template_path)
        
        # Generate slides
        for i, slide_data in enumerate(slides_content):
            print(f"    生成第 {i+1}/{len(slides_content)} 页")
            
            layout_index = slide_data.get('layout_index', 1)
            
            # Ensure layout index is valid
            if layout_index >= len(prs.slide_layouts):
                layout_index = 1  # Fallback to default content layout
            
            layout = prs.slide_layouts[layout_index]
            slide = prs.slides.add_slide(layout)
            
            # Fill slide content
            self._fill_slide(slide, slide_data)
        
        # Save presentation
        output_path = 'output.pptx'
        prs.save(output_path)
        
        print(f"  ✓ 已生成: {output_path}")
        return output_path
    
    def _fill_slide(self, slide, slide_data: Dict):
        """Fill a slide with content"""
        # Fill title
        if slide.shapes.title and 'title' in slide_data:
            slide.shapes.title.text = slide_data['title']
        
        # Fill body content
        body_content = slide_data.get('body', [])
        
        # Find body placeholder
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 2:  # BODY placeholder
                self._fill_text_frame(shape.text_frame, body_content)
                break
        
        # Handle images if present
        if 'images' in slide_data:
            self._insert_images(slide, slide_data['images'])
        
        # Handle tables if present
        if 'table_data' in slide_data:
            self._insert_table(slide, slide_data['table_data'])
    
    def _fill_text_frame(self, text_frame, body_content):
        """Fill text frame with content"""
        text_frame.clear()
        
        if isinstance(body_content, list):
            # List of bullet points
            for i, item in enumerate(body_content):
                if i == 0:
                    text_frame.text = item
                else:
                    p = text_frame.add_paragraph()
                    p.text = item
                    p.level = 0
        elif isinstance(body_content, str):
            # Plain text
            text_frame.text = body_content
        
        # Auto-fit text if enabled
        if self.config.get('generation', {}).get('font_size_adjustment', True):
            self._auto_fit_text(text_frame)
    
    def _auto_fit_text(self, text_frame, max_font_size: int = 24, min_font_size: int = 12):
        """Auto-fit text by adjusting font size"""
        # Note: python-pptx doesn't provide direct overflow detection
        # This is a simplified implementation
        try:
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size and run.font.size > Pt(max_font_size):
                        run.font.size = Pt(max_font_size)
        except:
            pass
    
    def _insert_images(self, slide, images: List[Dict]):
        """Insert images into slide"""
        # This is a placeholder implementation
        # In a full implementation, you would:
        # 1. Find suitable positions for images
        # 2. Resize images appropriately
        # 3. Insert them into the slide
        pass
    
    def _insert_table(self, slide, table_data: List[List[str]]):
        """Insert table into slide"""
        if not table_data:
            return
        
        # Calculate table dimensions
        rows = len(table_data)
        cols = len(table_data[0]) if rows > 0 else 0
        
        if rows == 0 or cols == 0:
            return
        
        # Position and size (centered in slide)
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(0.5 * rows)
        
        # Add table
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Fill table data
        for i, row_data in enumerate(table_data):
            for j, cell_value in enumerate(row_data):
                if j < cols:  # Ensure we don't exceed column count
                    table.cell(i, j).text = str(cell_value)
